import os
import copy
from datetime import datetime
import xlsxwriter

def _sv(val):
    if val is None or str(val).strip() == "" or str(val).strip() == "-": return 0.0
    try: return float(val)
    except: return 0.0



# python-pptx imports
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    from io import BytesIO
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def export_excel_visual(data_dict, output_path, metadata):
    wb = xlsxwriter.Workbook(output_path)
    
    # Formats
    fmt_bg = wb.add_format({'bg_color': '#FFFFFF'})
    fmt_header = wb.add_format({'bg_color': '#1E3A8A', 'font_color': '#FFFFFF', 'bold': True, 'font_size': 14, 'align': 'left', 'valign': 'vcenter', 'indent': 1})
    fmt_header_sub = wb.add_format({'bg_color': '#1E3A8A', 'font_color': '#E0E7FF', 'font_size': 10, 'align': 'left', 'valign': 'top', 'indent': 1})
    fmt_sidebar = wb.add_format({'bg_color': '#F8FAFC', 'border': 1, 'border_color': '#E2E8F0'})
    fmt_sidebar_title = wb.add_format({'bg_color': '#1E40AF', 'font_color': '#FFFFFF', 'bold': True, 'align': 'center', 'valign': 'vcenter'})
    fmt_slicer_lbl = wb.add_format({'bold': True, 'font_size': 10, 'valign': 'bottom'})
    fmt_slicer_box = wb.add_format({'bg_color': '#FFFFFF', 'border': 1, 'border_color': '#CBD5E1', 'valign': 'vcenter', 'indent': 1})
    
    fmt_kpi_box = wb.add_format({'bg_color': '#FFFFFF', 'border': 1, 'border_color': '#E2E8F0'})
    fmt_kpi_val = wb.add_format({'bold': True, 'font_size': 20, 'font_color': '#1E293B', 'align': 'center', 'valign': 'vcenter'})
    fmt_kpi_val_green = wb.add_format({'bold': True, 'font_size': 18, 'font_color': '#16A34A', 'align': 'center', 'valign': 'vcenter'})
    fmt_kpi_lbl = wb.add_format({'font_size': 10, 'font_color': '#64748B', 'align': 'center', 'valign': 'top'})
    
    fmt_tbl_hdr = wb.add_format({'bg_color': '#DBEAFE', 'bold': True, 'font_color': '#1E3A8A', 'border': 1, 'border_color': '#BFDBFE'})
    fmt_tbl_row1 = wb.add_format({'bg_color': '#FFFFFF', 'border': 1, 'border_color': '#E2E8F0', 'num_format': '#,##0'})
    fmt_tbl_row2 = wb.add_format({'bg_color': '#F8FAFC', 'border': 1, 'border_color': '#E2E8F0', 'num_format': '#,##0'})
    fmt_tbl_total = wb.add_format({'bg_color': '#EFF6FF', 'bold': True, 'font_color': '#1E40AF', 'border': 1, 'border_color': '#BFDBFE', 'num_format': '#,##0'})

    # 1. Create Data Sheet (Hidden)
    ws_data = wb.add_worksheet("Data_Hidden")
    ws_data.hide()
    
    total_data = data_dict.get("Total AH Gunsar", {})
    periode_list = total_data.get("periode_list", [])
    
    # Write periods to row 0
    ws_data.write(0, 0, "Periode")
    for i, p in enumerate(periode_list):
        ws_data.write(0, i+1, p)
        
    # Write DPK and Pinjaman Total
    dk_vals = []
    pk_vals = []
    dk_total = 0
    pk_total = 0
    if periode_list:
        lp = periode_list[-1]
        for r in total_data.get("rows", []):
            ll = r.get("label", "").lower()
            if ll == "dana pihak ketiga":
                for i, p in enumerate(periode_list):
                    val = _sv(r.get("values", {}).get(p, 0))
                    dk_vals.append(val)
                    ws_data.write(1, i+1, val)
                dk_total = _sv(r.get("values", {}).get(lp, 0))
            elif ll == "pinjaman":
                for i, p in enumerate(periode_list):
                    val = _sv(r.get("values", {}).get(p, 0))
                    pk_vals.append(val)
                    ws_data.write(2, i+1, val)
                pk_total = _sv(r.get("values", {}).get(lp, 0))
                
    ws_data.write(1, 0, "Simpanan")
    ws_data.write(2, 0, "Pinjaman")
    
    # Write Top KC Data (Simpanan & Pinjaman) for latest period
    kc_list = ["Tanah Abang", "Krekot", "Veteran", "Roxi", "Gunung Sahari", "Mangga Dua", "Kemayoran"]
    kc_simpanan = []
    kc_pinjaman = []
    for r_idx, kc in enumerate(kc_list):
        s_val, p_val = 0, 0
        if kc in data_dict and periode_list:
            lp = periode_list[-1]
            kc_r = data_dict[kc].get("rows", [])
            for x in kc_r:
                ll = x.get("label", "").lower()
                if ll == "dana pihak ketiga": s_val = _sv(x.get("values", {}).get(lp, 0))
                elif ll == "pinjaman": p_val = _sv(x.get("values", {}).get(lp, 0))
        kc_simpanan.append((kc, s_val))
        kc_pinjaman.append((kc, p_val))
        
    # Sort for Top 5
    top_5_kc = sorted(kc_simpanan, key=lambda x: x[1], reverse=True)[:5]
    ws_data.write(4, 0, "Top KC")
    ws_data.write(4, 1, "Simpanan")
    for i, (k, v) in enumerate(top_5_kc):
        ws_data.write(5+i, 0, k)
        ws_data.write(5+i, 1, v)
        
    # Distribusi per KC (Simpanan)
    ws_data.write(12, 0, "KC Distribusi")
    ws_data.write(12, 1, "Total")
    for i, (k, v) in enumerate(kc_simpanan):
        ws_data.write(13+i, 0, k)
        ws_data.write(13+i, 1, v)

    # 2. Create Dashboard Sheet
    ws = wb.add_worksheet("Dashboard")
    ws.hide_gridlines(2)
    ws.set_zoom(80)
    
    # Setup columns (A: Sidebar, B-M: Grid)
    ws.set_column('A:A', 22, fmt_bg)
    ws.set_column('B:L', 12, fmt_bg)
    ws.set_column('M:M', 2, fmt_bg) # margin
    for i in range(1, 40):
        ws.set_row(i, 15, fmt_bg)
        
    # Header Row 1-3
    ws.merge_range('A1:L2', "Visualisasi Periode", fmt_header)
    ws.merge_range('A3:L3', f"Data per {metadata.get('tanggal', '')} {metadata.get('hari', '')}, {metadata.get('jam', '')} WIB", fmt_header_sub)
    ws.set_row(0, 24)
    
    # Sidebar (A5:A30)
    ws.write('A5', "SLICER / FILTER", fmt_sidebar_title)
    filters = ["Periode", "Tanggal", "KC", "Jenis Data", "Sumber Data"]
    r_f = 6
    for f in filters:
        ws.write(r_f, 0, f, fmt_slicer_lbl)
        ws.write(r_f+1, 0, "(Semua)", fmt_slicer_box)
        r_f += 3
        
    # KPI Cards (Row 5-7, cols B-L)
    ws.merge_range('B5:C5', "TOTAL FILE", fmt_kpi_lbl)
    ws.merge_range('B6:C7', "3 / 3", fmt_kpi_val)
    
    ws.merge_range('D5:F5', "TOTAL KC", fmt_kpi_lbl)
    ws.merge_range('D6:F7', "7", fmt_kpi_val)
    
    ws.merge_range('G5:I5', "PERIODE SSA", fmt_kpi_lbl)
    ws.merge_range('G6:I7', periode_list[-1] if periode_list else "", fmt_kpi_val)
    
    ws.merge_range('J5:L5', "STATUS DASHBOARD", fmt_kpi_lbl)
    ws.merge_range('J6:L7', "READY", fmt_kpi_val_green)
    
    # --- CHARTS ---
    
    # Chart 1: Trend Simpanan & Pinjaman (Col B-G, Row 9-19)
    chart1 = wb.add_chart({'type': 'column'})
    if periode_list:
        chart1.add_series({
            'name': 'Simpanan',
            'categories': ['Data_Hidden', 0, 1, 0, len(periode_list)],
            'values':     ['Data_Hidden', 1, 1, 1, len(periode_list)],
            'fill':       {'color': '#004C97'},
            'data_labels': {'value': True, 'num_format': '#,##0'}
        })
        chart1.add_series({
            'name': 'Pinjaman',
            'categories': ['Data_Hidden', 0, 1, 0, len(periode_list)],
            'values':     ['Data_Hidden', 2, 1, 2, len(periode_list)],
            'fill':       {'color': '#F97316'},
            'data_labels': {'value': True, 'num_format': '#,##0'}
        })
    chart1.set_title({'name': 'TREND SIMPANAN & PINJAMAN (Rp Juta)', 'name_font': {'size': 11, 'bold': True}})
    chart1.set_legend({'position': 'top'})
    chart1.set_chartarea({'border': {'color': '#E2E8F0'}})
    ws.insert_chart('B9', chart1, {'x_scale': 1.1, 'y_scale': 1.4})
    
    # Chart 2: Top 5 KC (Col H-L, Row 9-19)
    chart2 = wb.add_chart({'type': 'bar'})
    chart2.add_series({
        'name': 'Total Portofolio',
        'categories': ['Data_Hidden', 5, 0, 9, 0],
        'values':     ['Data_Hidden', 5, 1, 9, 1],
        'fill':       {'color': '#2563EB'},
        'data_labels': {'value': True, 'num_format': '#,##0'}
    })
    chart2.set_title({'name': 'TOP 5 KC - TOTAL SIMPANAN (Rp Juta)', 'name_font': {'size': 11, 'bold': True}})
    chart2.set_legend({'none': True})
    chart2.set_chartarea({'border': {'color': '#E2E8F0'}})
    ws.insert_chart('H9', chart2, {'x_scale': 1.1, 'y_scale': 1.4})
    
    # Chart 3: Distribusi KC (Col B-D, Row 21-30)
    chart3 = wb.add_chart({'type': 'doughnut'})
    chart3.add_series({
        'name': 'Distribusi KC',
        'categories': ['Data_Hidden', 13, 0, 19, 0],
        'values':     ['Data_Hidden', 13, 1, 19, 1],
        'data_labels': {'percentage': True}
    })
    chart3.set_title({'name': 'DISTRIBUSI PORTOFOLIO PER KC', 'name_font': {'size': 11, 'bold': True}})
    chart3.set_chartarea({'border': {'color': '#E2E8F0'}})
    ws.insert_chart('B21', chart3, {'x_scale': 0.8, 'y_scale': 1.2})
    
    # Chart 4: Perkembangan Line (Col E-G, Row 21-30)
    chart4 = wb.add_chart({'type': 'line'})
    if periode_list:
        chart4.add_series({
            'name': 'Simpanan',
            'categories': ['Data_Hidden', 0, 1, 0, len(periode_list)],
            'values':     ['Data_Hidden', 1, 1, 1, len(periode_list)],
            'line':       {'color': '#004C97', 'width': 2.25},
            'marker':     {'type': 'circle', 'size': 6, 'border': {'color': '#004C97'}, 'fill': {'color': '#FFFFFF'}}
        })
        chart4.add_series({
            'name': 'Pinjaman',
            'categories': ['Data_Hidden', 0, 1, 0, len(periode_list)],
            'values':     ['Data_Hidden', 2, 1, 2, len(periode_list)],
            'line':       {'color': '#F97316', 'width': 2.25},
            'marker':     {'type': 'circle', 'size': 6, 'border': {'color': '#F97316'}, 'fill': {'color': '#FFFFFF'}}
        })
    chart4.set_title({'name': 'PERKEMBANGAN DATA (Rp Juta)', 'name_font': {'size': 11, 'bold': True}})
    chart4.set_legend({'position': 'top'})
    chart4.set_chartarea({'border': {'color': '#E2E8F0'}})
    ws.insert_chart('F21', chart4, {'x_scale': 1.0, 'y_scale': 1.2})
    
    # Table: Komposisi Data (Col J-L, Row 21-28)
    r_t = 20
    ws.write(r_t, 9, "Komponen", fmt_tbl_hdr)
    ws.write(r_t, 10, "Simpanan", fmt_tbl_hdr)
    ws.write(r_t, 11, "Pinjaman", fmt_tbl_hdr)
    
    r_t += 1
    components = ["Tabungan", "Giro", "Deposito"]
    total_rows = total_data.get("rows", [])
    for comp in components:
        s_val = 0
        if periode_list:
            lp = periode_list[-1]
            for r in total_rows:
                lbl = r.get("label", "").strip()
                if lbl.lower() == comp.lower():
                    s_val = _sv(r.get("values", {}).get(lp, 0))
        fmt = fmt_tbl_row1 if r_t % 2 == 0 else fmt_tbl_row2
        ws.write(r_t, 9, comp, fmt)
        ws.write(r_t, 10, s_val, fmt)
        ws.write(r_t, 11, "-", fmt)
        r_t += 1
        
    ws.write(r_t, 9, "TOTAL", fmt_tbl_total)
    ws.write(r_t, 10, dk_total, fmt_tbl_total)
    ws.write(r_t, 11, pk_total, fmt_tbl_total)
    
    wb.close()


def _generate_html(data_dict, metadata):
    # Buat HTML untuk Weasyprint
    periode = metadata.get("periode", "-")
    tanggal = f"{metadata.get('tanggal', '')} {metadata.get('hari', '')}, {metadata.get('jam', '')} WIB"
    
    html = f"""
    <html>
    <head>
    <style>
        @page {{ size: A4 landscape; margin: 20mm; }}
        body {{ font-family: Helvetica, Arial, sans-serif; color: #0F172A; font-size: 11px; }}
        .cover {{ background-color: #1E3A5F; color: white; height: 100vh; text-align: center; display: flex; flex-direction: column; justify-content: center; position: relative; page-break-after: always; }}
        .cover-bottom {{ border-bottom: 8px solid #F97316; position: absolute; bottom: 0; width: 100%; }}
        .title {{ font-size: 36px; font-weight: bold; margin-bottom: 10px; }}
        .subtitle {{ font-size: 20px; margin-bottom: 30px; }}
        
        .page {{ page-break-after: always; }}
        .header {{ background-color: #1E3A5F; color: white; padding: 10px; font-size: 16px; font-weight: bold; margin-bottom: 20px; }}
        
        .kpi-container {{ display: flex; justify-content: space-between; margin-bottom: 20px; }}
        .kpi-box {{ width: 23%; padding: 15px; border-top: 4px solid #2563EB; background: #F8FAFC; text-align: center; border-radius: 4px; box-shadow: 0 1px 3px rgba(0,0,0,0.1); }}
        .kpi-val {{ font-size: 18px; font-weight: bold; color: #1E293B; }}
        .kpi-lbl {{ font-size: 10px; color: #64748B; }}
        
        table {{ width: 100%; border-collapse: collapse; margin-bottom: 20px; font-size: 9px; }}
        th, td {{ padding: 6px; border: 1px solid #CBD5E1; text-align: right; }}
        th {{ background-color: #1E3A5F; color: white; text-align: center; font-weight: bold; }}
        td.left {{ text-align: left; }}
        
        .sep {{ background-color: #1E3A5F; height: 4px; padding: 0; }}
        .bg-dpk {{ background-color: #DBEAFE; }}
        .bg-sml {{ background-color: #FEF3C7; }}
        .bg-npl {{ background-color: #FEF2F2; }}
        .bg-alt {{ background-color: #F8FAFC; }}
        
        .status-baik {{ background-color: #DCFCE7; color: #166534; font-weight: bold; text-align: center; }}
        .status-perhatian {{ background-color: #FEF9C3; color: #854D0E; font-weight: bold; text-align: center; }}
        .status-kritis {{ background-color: #FEE2E2; color: #991B1B; font-weight: bold; text-align: center; }}
        
        .footer {{ position: fixed; bottom: 0; width: 100%; text-align: right; font-size: 8px; color: #64748B; padding-top: 10px; border-top: 1px solid #E2E8F0; }}
    </style>
    </head>
    <body>
    """
    
    # 1. Cover
    html += f"""
    <div class="cover">
        <h1 class="title">DASHBOARD SSA</h1>
        <div class="subtitle">AH Gunsar Jakarta Region</div>
        <div>Periode Data: {periode}</div>
        <div>Diekspor: {tanggal}</div>
        <div class="cover-bottom"></div>
    </div>
    """
    
    # Get total data
    total_data = data_dict.get("Total AH Gunsar", {})
    periode_list = total_data.get("periode_list", [])
    
    dk, pk, sml, npl = 0, 0, 0, 0
    if periode_list:
        lp = periode_list[-1]
        for r in total_data.get("rows", []):
            ll = r.get("label", "").lower()
            vl = r.get("values", {}).get(lp, 0)
            if ll == "dana pihak ketiga": dk = vl
            elif ll == "pinjaman": pk = vl
            elif ll == "sml %": sml = vl
            elif ll == "npl %": npl = vl
            
    # 2. Ringkasan
    html += f"""
    <div class="page">
        <div class="header">RINGKASAN EKSEKUTIF - AH GUNSAR</div>
        <div class="kpi-container">
            <div class="kpi-box"><div class="kpi-val">{dk:,.0f}</div><div class="kpi-lbl">Total DPK (Juta)</div></div>
            <div class="kpi-box"><div class="kpi-val">{pk:,.0f}</div><div class="kpi-lbl">Total Pinjaman (Juta)</div></div>
            <div class="kpi-box"><div class="kpi-val">{sml:.2f}%</div><div class="kpi-lbl">SML Ratio</div></div>
            <div class="kpi-box"><div class="kpi-val">{npl:.2f}%</div><div class="kpi-lbl">NPL Ratio</div></div>
        </div>
        <table>
            <tr>
                <th>Kantor Cabang</th><th>DPK Total</th><th>Pinjaman</th><th>SML %</th><th>NPL %</th><th>Status</th>
            </tr>
    """
    
    for kc_nm in ["Tanah Abang", "Krekot", "Veteran", "Roxi", "Gunung Sahari", "Mangga Dua", "Kemayoran", "Total AH Gunsar"]:
        if kc_nm not in data_dict: continue
        kc_r = data_dict[kc_nm].get("rows", [])
        
        k_dk, k_pk, k_sml, k_npl = 0, 0, 0, 0
        if periode_list:
            lp = periode_list[-1]
            for x in kc_r:
                ll = x.get("label","").lower()
                vl = x.get("values", {}).get(lp, 0)
                if ll == "dana pihak ketiga": k_dk = vl
                elif ll == "pinjaman": k_pk = vl
                elif ll == "sml %": k_sml = vl
                elif ll == "npl %": k_npl = vl
        
        if k_npl < 5: status_cls = "status-baik"; status_txt = "BAIK"
        elif k_npl <= 8: status_cls = "status-perhatian"; status_txt = "PERHATIAN"
        else: status_cls = "status-kritis"; status_txt = "KRITIS"
        
        bold_st = "font-weight: bold;" if kc_nm == "Total AH Gunsar" else ""
        
        html += f"""
        <tr style="{bold_st}">
            <td class="left">{kc_nm}</td>
            <td>{k_dk:,.0f}</td>
            <td>{k_pk:,.0f}</td>
            <td>{k_sml:.2f}%</td>
            <td>{k_npl:.2f}%</td>
            <td class="{status_cls}">{status_txt}</td>
        </tr>
        """
    html += "</table></div>"
    
    # 3. Data Per KC
    for kc_nm in ["Tanah Abang", "Krekot", "Veteran", "Roxi", "Gunung Sahari", "Mangga Dua", "Kemayoran", "Total AH Gunsar"]:
        if kc_nm not in data_dict: continue
        html += f"""
        <div class="page">
            <div class="header">DETAIL DATA — {kc_nm}</div>
            <table>
                <tr>
                    <th class="left">Mata Anggaran</th>
        """
        for p in periode_list:
            html += f"<th>{p}</th>"
        html += "<th>Growth MTD</th></tr>"
        
        rows = data_dict[kc_nm].get("rows", [])
        for i, r in enumerate(rows):
            lbl = r.get("label", "")
            if lbl == "SEP":
                html += '<tr><td colspan="{}" class="sep"></td></tr>'.format(len(periode_list)+2)
                continue
                
            lbl_l = lbl.lower()
            lvl = r.get("level", 0)
            is_bold = r.get("is_bold", False)
            
            cls = ""
            if "dana pihak" in lbl_l or "dpk" in lbl_l or "pinjaman" in lbl_l:
                if is_bold: cls = "bg-dpk"
            elif "sml" in lbl_l: cls = "bg-sml"
            elif "npl" in lbl_l: cls = "bg-npl"
            elif lvl > 0 and i % 2 == 0: cls = "bg-alt"
            
            bld = "font-weight: bold;" if is_bold else ""
            ind = "&nbsp;" * (lvl * 4)
            
            html += f'<tr class="{cls}" style="{bld}"><td class="left">{ind}{lbl}</td>'
            for p in periode_list:
                v = r.get("values", {}).get(p, 0)
                if "%" in lbl: html += f"<td>{v:.2f}%</td>"
                else: html += f"<td>{v:,.0f}</td>"
                
            g = r.get("growth_mtd", 0)
            c = "color: #16A34A;" if g > 0 else "color: #DC2626;" if g < 0 else ""
            html += f'<td style="{c}">{g:.2f}%</td></tr>'
            
        html += "</table></div>"
        
    # Disclaimer
    html += f"""
    <div style="margin-top: 100px; text-align: center; color: #94A3B8; font-size: 14px;">
        Data bersifat rahasia dan hanya untuk penggunaan internal Bank BRI.<br/>
        Generated by SSA Dashboard.
    </div>
    </body></html>
    """
    return html


def export_pdf_visual(data_dict, output_path, metadata):
    # Using ReportLab to avoid Weasyprint GTK/Pango C-library issues on macOS
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors

    doc = SimpleDocTemplate(output_path, pagesize=landscape(A4),
                            rightMargin=30, leftMargin=30, topMargin=30, bottomMargin=30)
    elements = []
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#1E3A5F'),
        alignment=1,
        spaceAfter=10
    )
    subtitle_style = ParagraphStyle(
        'CustomSubtitle',
        parent=styles['Normal'],
        fontSize=14,
        textColor=colors.HexColor('#64748B'),
        alignment=1,
        spaceAfter=30
    )
    header_style = ParagraphStyle(
        'CustomHeader',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#1E3A5F'),
        spaceAfter=15
    )
    
    periode = metadata.get("periode", "-")
    tanggal = f"{metadata.get('tanggal', '')} {metadata.get('hari', '')}, {metadata.get('jam', '')} WIB"
    
    elements.append(Paragraph("DASHBOARD SSA — BANK BRI", title_style))
    elements.append(Paragraph("AH Gunsar Jakarta Region", subtitle_style))
    elements.append(Paragraph(f"Periode Data: {periode} | Diekspor: {tanggal}", subtitle_style))
    elements.append(Spacer(1, 40))
    
    # Get total data
    total_data = data_dict.get("Total AH Gunsar", {})
    periode_list = total_data.get("periode_list", [])
    
    dk, pk, sml, npl = 0, 0, 0, 0
    if periode_list:
        lp = periode_list[-1]
        for r in total_data.get("rows", []):
            ll = r.get("label", "").lower()
            vl = r.get("values", {}).get(lp, 0)
            if ll == "dana pihak ketiga": dk = _sv(vl)
            elif ll == "pinjaman": pk = _sv(vl)
            elif ll == "sml %": sml = _sv(vl)
            elif ll == "npl %": npl = _sv(vl)
            
    # KPI Table
    elements.append(Paragraph("RINGKASAN EKSEKUTIF", header_style))
    kpi_data = [
        ["Total DPK (Juta)", "Total Pinjaman (Juta)", "SML Ratio (%)", "NPL Ratio (%)"],
        [f"{dk:,.0f}", f"{pk:,.0f}", f"{sml:.2f}%", f"{npl:.2f}%"]
    ]
    kpi_table = Table(kpi_data, colWidths=[150]*4)
    kpi_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1E3A5F')),
        ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,0), 10),
        ('BOTTOMPADDING', (0,0), (-1,0), 8),
        ('BACKGROUND', (0,1), (-1,1), colors.HexColor('#F8FAFC')),
        ('TEXTCOLOR', (0,1), (-1,1), colors.HexColor('#1E293B')),
        ('FONTNAME', (0,1), (-1,1), 'Helvetica-Bold'),
        ('FONTSIZE', (0,1), (-1,1), 14),
        ('GRID', (0,0), (-1,-1), 1, colors.HexColor('#CBD5E1'))
    ]))
    elements.append(kpi_table)
    elements.append(Spacer(1, 20))
    
    # Summary Table per KC
    kcs = ["Tanah Abang", "Krekot", "Veteran", "Roxi", "Gunung Sahari", "Mangga Dua", "Kemayoran", "Total AH Gunsar"]
    sum_data = [["Kantor Cabang", "DPK Total", "Pinjaman Total", "SML %", "NPL %", "Status"]]
    
    for kc_nm in kcs:
        if kc_nm not in data_dict: continue
        kc_r = data_dict[kc_nm].get("rows", [])
        
        k_dk, k_pk, k_sml, k_npl = 0, 0, 0, 0
        if periode_list:
            lp = periode_list[-1]
            for x in kc_r:
                ll = x.get("label","").lower()
                vl = x.get("values", {}).get(lp, 0)
                if ll == "dana pihak ketiga": k_dk = _sv(vl)
                elif ll == "pinjaman": k_pk = _sv(vl)
                elif ll == "sml %": k_sml = _sv(vl)
                elif ll == "npl %": k_npl = _sv(vl)
                
        if k_npl < 5: status_txt = "BAIK"
        elif k_npl <= 8: status_txt = "PERHATIAN"
        else: status_txt = "KRITIS"
        
        sum_data.append([
            kc_nm, f"{k_dk:,.0f}", f"{k_pk:,.0f}", f"{k_sml:.2f}%", f"{k_npl:.2f}%", status_txt
        ])
        
    sum_table = Table(sum_data, colWidths=[150, 100, 100, 80, 80, 100])
    ts = TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1E3A5F')),
        ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
        ('ALIGN', (0,0), (-1,-1), 'RIGHT'),
        ('ALIGN', (0,0), (0,-1), 'LEFT'),
        ('ALIGN', (-1,0), (-1,-1), 'CENTER'),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('GRID', (0,0), (-1,-1), 1, colors.HexColor('#CBD5E1')),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#F8FAFC')])
    ])
    
    for i, row in enumerate(sum_data[1:], 1):
        if row[0] == "Total AH Gunsar":
            ts.add('FONTNAME', (0,i), (-1,i), 'Helvetica-Bold')
        if row[5] == "BAIK": ts.add('TEXTCOLOR', (5,i), (5,i), colors.HexColor('#166534'))
        elif row[5] == "PERHATIAN": ts.add('TEXTCOLOR', (5,i), (5,i), colors.HexColor('#854D0E'))
        else: ts.add('TEXTCOLOR', (5,i), (5,i), colors.HexColor('#991B1B'))
            
    sum_table.setStyle(ts)
    elements.append(sum_table)
    
    # Detail Pages per KC
    for kc_nm in kcs:
        if kc_nm not in data_dict: continue
        elements.append(PageBreak())
        elements.append(Paragraph(f"DETAIL DATA — {kc_nm.upper()}", header_style))
        
        headers = ["Mata Anggaran"] + periode_list + ["Growth MTD"]
        det_data = [headers]
        
        rows = data_dict[kc_nm].get("rows", [])
        
        ts_detail = TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1E3A5F')),
            ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
            ('ALIGN', (0,0), (-1,-1), 'RIGHT'),
            ('ALIGN', (0,0), (0,-1), 'LEFT'),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('FONTSIZE', (0,0), (-1,-1), 8),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
            ('LEFTPADDING', (0,0), (-1,-1), 4),
            ('RIGHTPADDING', (0,0), (-1,-1), 4),
        ])
        
        row_idx = 1
        for r in rows:
            lbl = r.get("label", "")
            if lbl == "SEP":
                det_data.append([""] * len(headers))
                ts_detail.add('BACKGROUND', (0,row_idx), (-1,row_idx), colors.HexColor('#1E3A5F'))
                row_idx += 1
                continue
                
            lvl = r.get("level", 0)
            is_bold = r.get("is_bold", False)
            
            lbl_l = lbl.lower()
            if "dana pihak" in lbl_l or "dpk" in lbl_l or "pinjaman" in lbl_l:
                if is_bold: ts_detail.add('BACKGROUND', (0,row_idx), (-1,row_idx), colors.HexColor('#DBEAFE'))
            elif "sml" in lbl_l: ts_detail.add('BACKGROUND', (0,row_idx), (-1,row_idx), colors.HexColor('#FEF3C7'))
            elif "npl" in lbl_l: ts_detail.add('BACKGROUND', (0,row_idx), (-1,row_idx), colors.HexColor('#FEF2F2'))
            elif lvl > 0 and row_idx % 2 == 0: ts_detail.add('BACKGROUND', (0,row_idx), (-1,row_idx), colors.HexColor('#F8FAFC'))
            
            if is_bold: ts_detail.add('FONTNAME', (0,row_idx), (-1,row_idx), 'Helvetica-Bold')
            
            ind = "  " * lvl
            row_data = [ind + lbl]
            for p in periode_list:
                v = _sv(r.get("values", {}).get(p, 0))
                if "%" in lbl: row_data.append(f"{v:.2f}%")
                else: row_data.append(f"{v:,.0f}")
                
            g = _sv(r.get("growth_mtd", 0))
            if g > 0: ts_detail.add('TEXTCOLOR', (-1,row_idx), (-1,row_idx), colors.HexColor('#16A34A'))
            elif g < 0: ts_detail.add('TEXTCOLOR', (-1,row_idx), (-1,row_idx), colors.HexColor('#DC2626'))
            row_data.append(f"{g:.2f}%")
            
            det_data.append(row_data)
            row_idx += 1
            
        # Distribute widths (Mata Anggaran gets more)
        total_w = 750
        cols_count = len(headers)
        w_main = 250
        w_other = (total_w - w_main) / max(1, cols_count - 1)
        
        det_table = Table(det_data, colWidths=[w_main] + [w_other]*(cols_count-1))
        det_table.setStyle(ts_detail)
        elements.append(det_table)
        
    doc.build(elements)


def export_pptx_visual(data_dict, output_path, metadata):
    if not PPTX_AVAILABLE:
        raise Exception("python-pptx library not available. Please install python-pptx.")
        
    prs = Presentation()
    # Set to widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 1. COVER SLIDE
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 58, 95) # #1E3A5F
    shape.line.color.rgb = RGBColor(30, 58, 95)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.33), Inches(1.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "DASHBOARD SSA — BANK BRI"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "AH Gunsar Jakarta Region"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf.add_paragraph()
    p3.text = f"Periode: {metadata.get('periode', '-')} | {metadata.get('tanggal', '')}"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(200, 200, 200)
    p3.alignment = PP_ALIGN.CENTER
    
    # Orange line
    shape = slide.shapes.add_shape(1, 0, Inches(7.3), prs.slide_width, Inches(0.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(249, 115, 22) # #F97316
    shape.line.fill.background()
    
    
    # 2. RINGKASAN EKSEKUTIF
    slide = prs.slides.add_slide(blank_slide_layout)
    # Header
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 58, 95)
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.text = "RINGKASAN EKSEKUTIF"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Extract total data
    total_data = data_dict.get("Total AH Gunsar", {})
    periode_list = total_data.get("periode_list", [])
    dk, pk, sml, npl = 0, 0, 0, 0
    if periode_list:
        lp = periode_list[-1]
        for r in total_data.get("rows", []):
            ll = r.get("label", "").lower()
            vl = r.get("values", {}).get(lp, 0)
            if ll == "dana pihak ketiga": dk = _sv(vl)
            elif ll == "pinjaman": pk = _sv(vl)
            elif ll == "sml %": sml = _sv(vl)
            elif ll == "npl %": npl = _sv(vl)
            
    # Draw KPI boxes
    box_w = Inches(3)
    box_h = Inches(1.2)
    top = Inches(1.2)
    labels = ["Total DPK (Juta)", "Total Pinjaman (Juta)", "SML Ratio (%)", "NPL Ratio (%)"]
    vals = [f"{dk:,.0f}", f"{pk:,.0f}", f"{sml:.2f}%", f"{npl:.2f}%"]
    
    for i in range(4):
        left = Inches(0.5 + i * 3.1)
        bx = slide.shapes.add_shape(1, left, top, box_w, box_h)
        bx.fill.solid(); bx.fill.fore_color.rgb = RGBColor(241, 245, 249)
        bx.line.color.rgb = RGBColor(37, 99, 235); bx.line.width = Pt(2)
        
        tf = bx.text_frame
        p = tf.add_paragraph()
        p.text = vals[i]
        p.font.bold = True; p.font.size = Pt(24); p.font.color.rgb = RGBColor(30, 41, 59)
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = labels[i]
        p.font.size = Pt(12); p.font.color.rgb = RGBColor(100, 116, 139)
        p.alignment = PP_ALIGN.CENTER
        
    # Table
    kcs = ["Tanah Abang", "Krekot", "Veteran", "Roxi", "Gunung Sahari", "Mangga Dua", "Kemayoran", "Total AH Gunsar"]
    rows = len(kcs) + 1
    cols = 5
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.8), Inches(12.33), Inches(4)).table
    
    tbl_headers = ["Kantor Cabang", "DPK Total", "Pinjaman Total", "SML %", "NPL %"]
    for i, h in enumerate(tbl_headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(30, 58, 95)
        for par in cell.text_frame.paragraphs:
            par.font.bold = True; par.font.color.rgb = RGBColor(255,255,255); par.font.size = Pt(14)
            
    for r_idx, kc_nm in enumerate(kcs):
        kc_r = data_dict.get(kc_nm, {}).get("rows", [])
        k_dk, k_pk, k_sml, k_npl = 0, 0, 0, 0
        if periode_list:
            lp = periode_list[-1]
            for x in kc_r:
                ll = x.get("label","").lower()
                vl = x.get("values", {}).get(lp, 0)
                if ll == "dana pihak ketiga": k_dk = _sv(vl)
                elif ll == "pinjaman": k_pk = _sv(vl)
                elif ll == "sml %": k_sml = _sv(vl)
                elif ll == "npl %": k_npl = _sv(vl)
                
        table.cell(r_idx+1, 0).text = kc_nm
        table.cell(r_idx+1, 1).text = f"{k_dk:,.0f}"
        table.cell(r_idx+1, 2).text = f"{k_pk:,.0f}"
        table.cell(r_idx+1, 3).text = f"{k_sml:.2f}%"
        table.cell(r_idx+1, 4).text = f"{k_npl:.2f}%"
        
        for c_idx in range(5):
            if (r_idx+1) % 2 == 0:
                table.cell(r_idx+1, c_idx).fill.solid()
                table.cell(r_idx+1, c_idx).fill.fore_color.rgb = RGBColor(241, 245, 249)
            if kc_nm == "Total AH Gunsar":
                for par in table.cell(r_idx+1, c_idx).text_frame.paragraphs:
                    par.font.bold = True
                
    # Function to create chart
    def add_chart_slide(title_text, fig):
        sl = prs.slides.add_slide(blank_slide_layout)
        sh = sl.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(30, 58, 95)
        tf = sh.text_frame
        p = tf.add_paragraph()
        p.text = title_text
        p.font.bold = True; p.font.size = Pt(24); p.font.color.rgb = RGBColor(255, 255, 255)
        
        buf = BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
        buf.seek(0)
        plt.close(fig)
        sl.shapes.add_picture(buf, Inches(1), Inches(1.5), Inches(11.33), Inches(5.5))

    
    # CHARTS using matplotlib
    plt.style.use('default')
    
    # Slide 3: Trend DPK & Pinjaman
    if periode_list and total_data:
        fig, ax = plt.subplots(figsize=(10, 5))
        dpks = []
        pjs = []
        for p in periode_list:
            d, p_val = 0, 0
            for r in total_data.get("rows", []):
                l = r.get("label", "").lower()
                v = r.get("values", {}).get(p, 0)
                if l == "dana pihak ketiga": d = v
                elif l == "pinjaman": p_val = v
            dpks.append(d)
            pjs.append(p_val)
            
        x = range(len(periode_list))
        w = 0.35
        ax.bar([i - w/2 for i in x], dpks, w, label='DPK', color='#2563EB')
        ax.bar([i + w/2 for i in x], pjs, w, label='Pinjaman', color='#F59E0B')
        ax.set_xticks(x)
        ax.set_xticklabels(periode_list)
        ax.legend()
        ax.grid(True, axis='y', alpha=0.3)
        add_chart_slide("Trend Simpanan & Pinjaman (Rp Juta)", fig)
        
    # Detail slides per KC
    for kc_nm in kcs:
        kc_r = data_dict.get(kc_nm, {}).get("rows", [])
        if not kc_r: continue
        
        sl = prs.slides.add_slide(blank_slide_layout)
        sh = sl.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(30, 58, 95)
        tf = sh.text_frame
        p = tf.add_paragraph()
        p.text = f"DETAIL KINERJA — {kc_nm.upper()}"
        p.font.bold = True; p.font.size = Pt(24); p.font.color.rgb = RGBColor(255, 255, 255)
        
        k_dk, k_pk, k_sml, k_npl = 0, 0, 0, 0
        if periode_list:
            lp = periode_list[-1]
            for x in kc_r:
                ll = x.get("label","").lower()
                vl = x.get("values", {}).get(lp, 0)
                if ll == "dana pihak ketiga": k_dk = _sv(vl)
                elif ll == "pinjaman": k_pk = _sv(vl)
                elif ll == "sml %": k_sml = _sv(vl)
                elif ll == "npl %": k_npl = _sv(vl)
                
        # Mini cards left
        y_pos = 1.5
        for t, v in [("Total DPK", f"{k_dk:,.0f}"), ("Pinjaman", f"{k_pk:,.0f}"), ("SML Ratio", f"{k_sml:.2f}%"), ("NPL Ratio", f"{k_npl:.2f}%")]:
            bx = sl.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(2.5), Inches(1))
            bx.fill.solid(); bx.fill.fore_color.rgb = RGBColor(248, 250, 252)
            bx.line.color.rgb = RGBColor(203, 213, 225)
            
            tf = bx.text_frame
            p = tf.add_paragraph()
            p.text = t
            p.font.size = Pt(12); p.font.color.rgb = RGBColor(100, 116, 139)
            
            p2 = tf.add_paragraph()
            p2.text = v
            p2.font.bold = True; p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(30, 58, 95)
            
            y_pos += 1.3
            
        # Table right
        cols = 2 + len(periode_list)
        main_rows = [x for x in kc_r if x.get("level", 0) == 0 and x.get("label") != "SEP"]
        t = sl.shapes.add_table(len(main_rows)+1, cols, Inches(3.5), Inches(1.5), Inches(9.33), Inches(5)).table
        
        t.cell(0, 0).text = "Mata Anggaran"
        for i, pd_lbl in enumerate(periode_list):
            t.cell(0, 1+i).text = pd_lbl
        t.cell(0, cols-1).text = "MTD"
        
        for c_idx in range(cols):
            t.cell(0, c_idx).fill.solid(); t.cell(0, c_idx).fill.fore_color.rgb = RGBColor(30, 58, 95)
            for par in t.cell(0, c_idx).text_frame.paragraphs:
                par.font.color.rgb = RGBColor(255, 255, 255)
                
        for r_idx, row in enumerate(main_rows):
            lbl = row.get("label", "")
            t.cell(r_idx+1, 0).text = lbl
            for i, p in enumerate(periode_list):
                v = _sv(row.get("values", {}).get(p, 0))
                if "%" in lbl: t.cell(r_idx+1, 1+i).text = f"{v:.2f}%"
                else: t.cell(r_idx+1, 1+i).text = f"{v:,.0f}"
                
            g = _sv(row.get("growth_mtd", 0))
            t.cell(r_idx+1, cols-1).text = f"{g:.2f}%"
            
    # Penutup
    sl = prs.slides.add_slide(blank_slide_layout)
    sh = sl.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(30, 58, 95)
    
    txBox = sl.shapes.add_textbox(Inches(2), Inches(3), Inches(9.33), Inches(1.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "TERIMA KASIH"
    p.font.bold = True; p.font.size = Pt(50); p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    prs.save(output_path)
